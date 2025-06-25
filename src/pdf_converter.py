import json
import os
import sys

# Importa i moduli necessari con gestione degli errori
try:
    import fitz  # PyMuPDF
except ImportError:
    print("Errore: PyMuPDF non installato.")
    print("Installa con: pip install PyMuPDF")
    print("Oppure esegui lo script di setup: ./setup.sh")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Errore: python-pptx non installato.")
    print("Installa con: pip install python-pptx")
    print("Oppure esegui lo script di setup: ./setup.sh")
    sys.exit(1)


def load_config(config_path="config.json"):
    """Carica la configurazione dal file JSON"""
    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except FileNotFoundError:
        print(f"File di configurazione {config_path} non trovato!")
        return None
    except json.JSONDecodeError:
        print(f"Errore nel parsing del file di configurazione {config_path}!")
        return None


def pdf_to_ppt(pdf_path, output_path, config):
    """
    Converte PDF in PowerPoint usando PyMuPDF (fitz)
    """
    print("Conversione PDF in PowerPoint usando PyMuPDF...")

    # Apri il PDF
    pdf_document = fitz.open(pdf_path)

    # Crea una nuova presentazione PowerPoint
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Layout vuoto

    # Ottieni le impostazioni dalla configurazione
    matrix_scale = config['output_settings']['image_quality']['matrix_scale']

    # Ottieni le dimensioni della prima pagina per impostare le slide
    first_page = pdf_document.load_page(0)
    page_rect = first_page.rect
    pdf_width = page_rect.width
    pdf_height = page_rect.height
    pdf_aspect_ratio = pdf_width / pdf_height

    # Imposta le dimensioni della slide basate sulle proporzioni del PDF
    # Usa una larghezza standard di 10 pollici e calcola l'altezza proporzionale
    slide_width_inches = 10
    slide_height_inches = slide_width_inches / pdf_aspect_ratio

    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)

    print(f"PDF ratio: {pdf_aspect_ratio:.2f}")
    print(
        f"Slide dimensions: {slide_width_inches}x{slide_height_inches:.2f} inches")

    for page_num in range(len(pdf_document)):
        print(f"Elaborando pagina {page_num + 1}/{len(pdf_document)}...")

        # Converti la pagina PDF in immagine
        page = pdf_document.load_page(page_num)

        pix = page.get_pixmap(matrix=fitz.Matrix(matrix_scale, matrix_scale))
        img_data = pix.tobytes("png")

        # Salva temporaneamente l'immagine
        temp_img_path = f"temp_page_{page_num}.png"
        with open(temp_img_path, "wb") as f:
            f.write(img_data)

        # Aggiungi slide alla presentazione
        slide = prs.slides.add_slide(slide_layout)

        # L'immagine occupa tutta la slide mantenendo le proporzioni originali
        slide.shapes.add_picture(temp_img_path, Inches(0), Inches(0),
                                 width=Inches(slide_width_inches),
                                 height=Inches(slide_height_inches))

        # Rimuovi il file temporaneo
        os.remove(temp_img_path)

    # Salva il numero di pagine prima di chiudere il documento
    num_pages = len(pdf_document)

    # Chiudi il PDF
    pdf_document.close()

    # Salva la presentazione
    prs.save(output_path)
    print(f"Conversione completata: {output_path}")
    print(f"Slide create: {num_pages}")
    return True


def convert_pdf_to_ppt(pdf_path=None, output_path=None):
    """
    Funzione principale che gestisce la conversione
    """
    # Carica la configurazione
    config = load_config()
    if not config:
        return False

    # Usa i valori dalla configurazione se non specificati
    if not pdf_path:
        pdf_path = config['file_paths']['input_pdf']
    if not output_path:
        output_path = config['output_settings']['default_output_name']

    # Verifica che il file PDF esista
    if not os.path.exists(pdf_path):
        print(f"Errore: File PDF {pdf_path} non trovato!")
        return False

    # Crea la directory di output se non esiste
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)

    print(f"Conversione di: {pdf_path}")
    print(f"Output: {output_path}")

    # Esegui la conversione
    return pdf_to_ppt(pdf_path, output_path, config)


if __name__ == "__main__":
    if len(sys.argv) > 1:
        if sys.argv[1] == "--help":
            print("Uso:")
            print(
                "  python pdf_converter.py                    # Usa configurazione predefinita")
            print(
                "  python pdf_converter.py input.pdf output.pptx  # Specifica parametri")
        elif len(sys.argv) >= 3:
            pdf_file = sys.argv[1]
            output_file = sys.argv[2]
            convert_pdf_to_ppt(pdf_file, output_file)
        else:
            print("Parametri insufficienti. Usa --help per vedere l'uso.")
    else:
        convert_pdf_to_ppt()
