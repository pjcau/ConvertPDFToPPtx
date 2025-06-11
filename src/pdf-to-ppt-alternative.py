import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io
import os

def pdf_to_ppt(pdf_path, output_path):
    # Apri il PDF
    pdf_document = fitz.open(pdf_path)
    
    # Crea una nuova presentazione PowerPoint
    prs = Presentation()
    
    # Rimuovi la slide di default
    slide_layout = prs.slide_layouts[6]  # Layout vuoto
    
    for page_num in range(len(pdf_document)):
        # Converti la pagina PDF in immagine
        page = pdf_document.load_page(page_num)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Alta risoluzione
        img_data = pix.tobytes("png")
        
        # Salva temporaneamente l'immagine
        temp_img_path = f"temp_page_{page_num}.png"
        with open(temp_img_path, "wb") as f:
            f.write(img_data)
        
        # Aggiungi slide alla presentazione
        slide = prs.slides.add_slide(slide_layout)
        
        # Aggiungi l'immagine alla slide
        slide.shapes.add_picture(temp_img_path, Inches(0), Inches(0), 
                                width=Inches(10), height=Inches(7.5))
        
        # Rimuovi il file temporaneo
        os.remove(temp_img_path)
    
    # Chiudi il PDF
    pdf_document.close()
    
    # Salva la presentazione
    prs.save(output_path)
    print(f"Conversione completata: {output_path}")

if __name__ == "__main__":
    pdf_to_ppt("test.pdf", "pdf-to-ppt-free.pptx")